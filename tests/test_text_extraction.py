import io

import pytest
from docx import Document as DocxDocument
from pptx import Presentation

from services.text_extraction import (
    EMPTY_DOCUMENT_WARNING,
    SCANNED_PDF_WARNING,
    UnsupportedFileTypeError,
    extract_text,
)


def _build_minimal_pdf(text: str) -> bytes:
    """Construye un PDF válido mínimo a mano (sin dependencias extra) para pruebas."""
    content = f"BT /F1 24 Tf 100 700 Td ({text}) Tj ET".encode("latin-1") if text else b""
    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /Resources << /Font << /F1 4 0 R >> >> "
        b"/MediaBox [0 0 612 792] /Contents 5 0 R >>",
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
        b"<< /Length " + str(len(content)).encode() + b" >>\nstream\n" + content + b"\nendstream",
    ]
    out = bytearray(b"%PDF-1.4\n")
    offsets = [0]
    for index, body in enumerate(objects, start=1):
        offsets.append(len(out))
        out += f"{index} 0 obj\n".encode()
        out += body
        out += b"\nendobj\n"
    xref_offset = len(out)
    out += f"xref\n0 {len(objects) + 1}\n".encode()
    out += b"0000000000 65535 f \n"
    for offset in offsets[1:]:
        out += f"{offset:010d} 00000 n \n".encode()
    out += b"trailer\n"
    out += f"<< /Size {len(objects) + 1} /Root 1 0 R >>\n".encode()
    out += b"startxref\n"
    out += f"{xref_offset}\n".encode()
    out += b"%%EOF"
    return bytes(out)


def test_extract_txt():
    result = extract_text("apuntes.txt", io.BytesIO("Línea uno.\nLínea dos.".encode("utf-8")))
    assert "Línea uno." in result.text
    assert "Línea dos." in result.text
    assert result.warnings == []


def test_extract_markdown():
    result = extract_text("resumen.md", io.BytesIO(b"# Titulo\n\nContenido del resumen."))
    assert "Titulo" in result.text
    assert "Contenido del resumen." in result.text


def test_extract_pdf_with_text():
    pdf_bytes = _build_minimal_pdf("Hola examen de abogacia")
    result = extract_text("examen.pdf", io.BytesIO(pdf_bytes))
    assert result.page_count == 1
    assert "Hola examen de abogacia" in result.text
    assert result.warnings == []


def test_extract_scanned_pdf_without_text():
    pdf_bytes = _build_minimal_pdf("")
    result = extract_text("escaneado.pdf", io.BytesIO(pdf_bytes))
    assert result.text == ""
    assert SCANNED_PDF_WARNING in result.warnings


def test_extract_docx():
    doc = DocxDocument()
    doc.add_paragraph("Primer párrafo de apuntes.")
    doc.add_paragraph("Segundo párrafo con más contenido.")
    buffer = io.BytesIO()
    doc.save(buffer)
    buffer.seek(0)
    result = extract_text("apuntes.docx", buffer)
    assert "Primer párrafo de apuntes." in result.text
    assert "Segundo párrafo con más contenido." in result.text


def test_extract_pptx():
    presentation = Presentation()
    slide_layout = presentation.slide_layouts[6]  # en blanco
    slide = presentation.slides.add_slide(slide_layout)
    textbox = slide.shapes.add_textbox(0, 0, 100, 100)
    textbox.text_frame.text = "Diapositiva de ejemplo"
    buffer = io.BytesIO()
    presentation.save(buffer)
    buffer.seek(0)
    result = extract_text("presentacion.pptx", buffer)
    assert result.page_count == 1
    assert "Diapositiva de ejemplo" in result.text


def test_extract_empty_document():
    result = extract_text("vacio.txt", io.BytesIO(b""))
    assert result.text == ""
    assert EMPTY_DOCUMENT_WARNING in result.warnings


def test_extract_unsupported_extension():
    with pytest.raises(UnsupportedFileTypeError):
        extract_text("archivo.xyz", io.BytesIO(b"contenido"))
