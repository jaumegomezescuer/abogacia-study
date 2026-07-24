"""Extracción de texto de documentos (PDF, DOCX, PPTX, TXT, MD).

Los archivos llegan como bytes en memoria (no hay carpeta local persistente:
todo se guarda como BLOB en la base de datos), así que se trabaja siempre
sobre un `BytesIO`, nunca sobre una ruta en disco.
"""
from __future__ import annotations

import re
from dataclasses import dataclass, field
from pathlib import Path
from typing import BinaryIO

from docx import Document as DocxDocument
from pptx import Presentation
from pypdf import PdfReader

SUPPORTED_EXTENSIONS = (".pdf", ".docx", ".pptx", ".txt", ".md")

SCANNED_PDF_WARNING = "scanned_pdf"
EMPTY_DOCUMENT_WARNING = "empty_document"


class UnsupportedFileTypeError(Exception):
    pass


@dataclass
class ExtractedDocument:
    text: str
    page_count: int | None
    pages: list[dict] = field(default_factory=list)
    warnings: list[str] = field(default_factory=list)


def _clean_text(text: str) -> str:
    text = text.replace("\x00", "")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def _extract_pdf(file_obj: BinaryIO) -> ExtractedDocument:
    reader = PdfReader(file_obj)
    pages: list[dict] = []
    for index, page in enumerate(reader.pages, start=1):
        raw = page.extract_text() or ""
        pages.append({"page": index, "text": _clean_text(raw)})
    text = "\n\n".join(p["text"] for p in pages if p["text"])
    warnings = [] if text.strip() else [SCANNED_PDF_WARNING]
    return ExtractedDocument(text=text, page_count=len(reader.pages), pages=pages, warnings=warnings)


def _extract_docx(file_obj: BinaryIO) -> ExtractedDocument:
    doc = DocxDocument(file_obj)
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                if cell.text.strip():
                    parts.append(cell.text)
    text = _clean_text("\n".join(parts))
    warnings = [] if text.strip() else [EMPTY_DOCUMENT_WARNING]
    return ExtractedDocument(text=text, page_count=None, pages=[{"page": 1, "text": text}], warnings=warnings)


def _extract_pptx(file_obj: BinaryIO) -> ExtractedDocument:
    presentation = Presentation(file_obj)
    pages: list[dict] = []
    for index, slide in enumerate(presentation.slides, start=1):
        lines = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in paragraph.runs)
                    if line.strip():
                        lines.append(line)
        pages.append({"page": index, "text": _clean_text("\n".join(lines))})
    text = "\n\n".join(p["text"] for p in pages if p["text"])
    warnings = [] if text.strip() else [EMPTY_DOCUMENT_WARNING]
    return ExtractedDocument(text=text, page_count=len(presentation.slides), pages=pages, warnings=warnings)


def _extract_plain_text(file_obj: BinaryIO) -> ExtractedDocument:
    raw_bytes = file_obj.read()
    try:
        raw = raw_bytes.decode("utf-8")
    except UnicodeDecodeError:
        raw = raw_bytes.decode("latin-1", errors="replace")
    text = _clean_text(raw)
    warnings = [] if text.strip() else [EMPTY_DOCUMENT_WARNING]
    return ExtractedDocument(text=text, page_count=None, pages=[{"page": 1, "text": text}], warnings=warnings)


def extract_text(filename: str, file_obj: BinaryIO) -> ExtractedDocument:
    suffix = Path(filename).suffix.lower()
    file_obj.seek(0)
    if suffix == ".pdf":
        return _extract_pdf(file_obj)
    if suffix == ".docx":
        return _extract_docx(file_obj)
    if suffix == ".pptx":
        return _extract_pptx(file_obj)
    if suffix in (".txt", ".md"):
        return _extract_plain_text(file_obj)
    raise UnsupportedFileTypeError(f"Extensión no admitida: {suffix or '(ninguna)'}")
